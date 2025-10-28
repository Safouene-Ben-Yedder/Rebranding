# app.py
# Run with: streamlit run app.py

import io
import os
import zipfile
import xml.etree.ElementTree as ET
from typing import Dict, List, Optional, Set, Tuple, Any
import base64

import streamlit as st

def get_base64_image(image_path):
    with open("pwc.png", "rb") as img_file:
        return base64.b64encode(img_file.read()).decode()

# Convert local image to base64
img_base64 = get_base64_image("pwc.png")

st.markdown(f"""
    <style>
    .fixed-image {{
        position: fixed;
        top: 8px;
        left: 20px;
        z-index: 999;
        width: 100px;
    }}
    </style>
    <img src="data:image/png;base64,{img_base64}" class="fixed-image">
    """, unsafe_allow_html=True)
st.write("")
st.write("")
st.write("")
st.write("")



# Brand styling (PwC-inspired)
PWC_ORANGE = "#E87722"
PWC_DARK_GRAY = "#3C3C3C"
PWC_LIGHT_GRAY = "#F5F5F5"
PWC_WHITE = "#FFFFFF"

# Provide helpful notes for missing libraries
DOCX_AVAILABLE = False
PPTX_AVAILABLE = False
OPENPYXL_AVAILABLE = False
PIL_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    from docx.shared import RGBColor
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor as PPTX_RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.dml import MSO_FILL, MSO_THEME_COLOR
    from pptx.oxml.ns import qn as pptx_qn
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Color, Border, Side
    from openpyxl.drawing.image import Image as XLImage
    OPENPYXL_AVAILABLE = True
except Exception:
    OPENPYXL_AVAILABLE = False

try:
    from PIL import Image as PILImage
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False


# Page configuration
st.set_page_config(
    page_title="PwC Rebranding Tool",
    page_icon="🎨",
    layout="wide",
    menu_items={"Get help": None, "Report a bug": None, "About": None}
)

FOOTER_TEXT = "© EMEA My Way Technology Team"

# Custom styling
st.markdown(f"""
    <style>
        .stApp {{background-color: {PWC_LIGHT_GRAY};}}
        .block-container {{padding-top: 1.5rem;}}
        .pwc-header {{background: linear-gradient(90deg, {PWC_ORANGE}, #f14e00);color: {PWC_WHITE};padding: 16px 20px;border-radius: 8px;margin-bottom: 20px;}}
        .custom-footer {{position: fixed; left: 0; right: 0; bottom: 0; padding: 10px 16px; background: white; border-top: 3px solid {PWC_ORANGE}; text-align: center; color: #666;font-size: 0.9rem; z-index: 9999;}}
        .pwc-card {{background-color: {PWC_WHITE};border: 1px solid #eaeaea;padding: 16px;border-radius: 8px;margin-bottom: 16px;}}
        .pwc-section-title {{color: {PWC_DARK_GRAY};font-weight: 700;margin-bottom: 8px;}}
        .pwc-subtle {{color: #ffffff;font-size: 0.9rem;}}
        .stButton>button {{background-color: {PWC_ORANGE};color: {PWC_WHITE};border-radius: 6px;border: none;}}
        .stButton>button:hover {{background-color: #ffffff;}}
        .pwc-hint {{font-size: 0.9rem;color: #666;}}
        .thumb-row {{display: flex; gap: 8px; flex-wrap: wrap; align-items: center;}}
        .thumb-item {{border: 1px solid #eee; border-radius: 6px; padding: 4px; background: #fafafa;}}
        #MainMenu {{visibility: hidden;}}
        header {{visibility: hidden;}}
        footer {{visibility: hidden;}}
        [data-testid="stToolbar"] {{display: none !important;}}
        [data-testid="stDecoration"] {{display: none !important;}}
        [data-testid="baseButton-header"] {{display: none !important;}}
        [data-testid="collapsedControl"] {{display: none !important;}}
    </style>
""", unsafe_allow_html=True)


# Session state
if "image_repls" not in st.session_state:
    st.session_state["image_repls"] = {}
if "theme_image_repls" not in st.session_state:
    st.session_state["theme_image_repls"] = {}

def get_persisted_image_replacements() -> Tuple[Dict[str, bytes], Dict[str, bytes]]:
    return dict(st.session_state["image_repls"]), dict(st.session_state["theme_image_repls"])

def persist_image_replacement(uid: str, data: bytes):
    st.session_state["image_repls"][uid] = data

def persist_theme_image_replacement(media_path: str, data: bytes):
    st.session_state["theme_image_repls"][media_path] = data

def clear_all_replacements():
    st.session_state["image_repls"].clear()
    st.session_state["theme_image_repls"].clear()


# General utils
def infer_file_type(file_name: str) -> Optional[str]:
    name = file_name.lower()
    if name.endswith(".docx"):
        return "docx"
    if name.endswith(".pptx"):
        return "pptx"
    if name.endswith(".xlsx"):
        return "xlsx"
    if name.endswith(".pdf"):
        return "pdf"
    return None

def hex_no_hash(hex_color: str) -> str:
    if not hex_color:
        return ""
    return hex_color.replace("#", "").upper()

def rgbcolor_to_hex(rgb_obj) -> Optional[str]:
    try:
        if rgb_obj is None:
            return None
        s = str(rgb_obj)
        if s and len(s) == 6:
            return "#" + s.upper()
        if hasattr(rgb_obj, "rgb"):
            val = getattr(rgb_obj, "rgb")
            if isinstance(val, int):
                r = (val >> 16) & 0xFF
                g = (val >> 8) & 0xFF
                b = val & 0xFF
                return "#{:02X}{:02X}{:02X}".format(r, g, b)
    except Exception:
        return None
    return None

def extract_color_from_pptx_color_obj(color_obj) -> Optional[str]:
    """Enhanced color extraction for PPTX that handles theme colors and RGB"""
    try:
        if color_obj is None:
            return None
        
        # Try RGB color first
        try:
            if hasattr(color_obj, 'rgb') and color_obj.rgb is not None:
                return rgbcolor_to_hex(color_obj.rgb)
        except Exception:
            pass
        
        # Try theme color
        try:
            if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
                # Theme colors can't be directly converted, but we can note them
                # For now, try to get the RGB value if available
                pass
        except Exception:
            pass
        
        # Try accessing color through XML
        try:
            if hasattr(color_obj, '_color') and color_obj._color is not None:
                color_elem = color_obj._color
                # Check for srgbClr (RGB color)
                srgb = color_elem.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if srgb is not None:
                    val = srgb.get('val')
                    if val and len(val) == 6:
                        return "#" + val.upper()
        except Exception:
            pass
            
    except Exception:
        pass
    return None

def openpyxl_color_to_hex(color_obj) -> Optional[str]:
    try:
        if color_obj is None:
            return None
        if hasattr(color_obj, "rgb") and color_obj.rgb:
            val = color_obj.rgb
            if len(val) == 8:
                return "#" + val[2:].upper()
            if len(val) == 6:
                return "#" + val.upper()
    except Exception:
        return None
    return None

def safe_key(s: str) -> str:
    return "".join(ch if ch.isalnum() or ch in "._-:" else "_" for ch in str(s))

def normalize_zip_path(base_dir: str, rel_target: str) -> str:
    parts = base_dir.strip("/").split("/")
    target_parts = rel_target.split("/")
    stack = parts[:]
    for seg in target_parts:
        if seg == "..":
            if stack:
                stack.pop()
        elif seg == "." or seg == "":
            continue
        else:
            stack.append(seg)
    return "/".join(stack)

def convert_image_bytes_to_ext(data: bytes, target_ext: str) -> bytes:
    if not PIL_AVAILABLE:
        return data
    try:
        ext = target_ext.lower().replace(".", "")
        img = PILImage.open(io.BytesIO(data))
        if ext in ("png", "gif", "webp"):
            img = img.convert("RGBA")
        else:
            img = img.convert("RGB")
        out = io.BytesIO()
        save_fmt = "PNG" if ext == "png" else "JPEG" if ext in ("jpg", "jpeg") else "PNG"
        img.save(out, format=save_fmt)
        return out.getvalue()
    except Exception:
        return data


# DOCX functions
def docx_extract_deep_formatting(element, text_colors: Set[str], shape_colors: Set[str], fonts: Set[str]):
    """Recursively extract formatting from nested DOCX elements"""
    try:
        # Check runs within paragraphs
        runs = element.xpath('.//w:r', namespaces=element.nsmap if hasattr(element, 'nsmap') else {})
        for run in runs:
            # Font name
            font_elems = run.xpath('.//w:rFonts', namespaces=run.nsmap if hasattr(run, 'nsmap') else {})
            for font_elem in font_elems:
                for attr in ['w:ascii', 'w:hAnsi', 'w:eastAsia', 'w:cs']:
                    font_name = font_elem.get(qn(attr))
                    if font_name:
                        fonts.add(font_name)
            
            # Text colors
            color_elems = run.xpath('.//w:color', namespaces=run.nsmap if hasattr(run, 'nsmap') else {})
            for color_elem in color_elems:
                color_val = color_elem.get(qn('w:val'))
                if color_val and color_val != "auto" and len(color_val) >= 6:
                    text_colors.add("#" + color_val[:6].upper())
        
        # Check shading and fills
        shd_elems = element.xpath('.//w:shd', namespaces=element.nsmap if hasattr(element, 'nsmap') else {})
        for shd in shd_elems:
            fill_val = shd.get(qn('w:fill'))
            if fill_val and fill_val != "auto" and len(fill_val) >= 6:
                shape_colors.add("#" + fill_val[:6].upper())
        
        # Check borders
        border_elems = element.xpath('.//w:*[contains(local-name(), "Border")]/*', 
                                   namespaces=element.nsmap if hasattr(element, 'nsmap') else {})
        for border in border_elems:
            color_val = border.get(qn('w:color'))
            if color_val and color_val.lower() not in ("auto", "none") and len(color_val) >= 6:
                shape_colors.add("#" + color_val[:6].upper())
                
    except Exception:
        pass

def docx_extract(file_bytes: bytes):
    if not DOCX_AVAILABLE:
        return None
    buf = io.BytesIO(file_bytes)
    doc = DocxDocument(buf)
    text_colors: Set[str] = set()
    fonts: Set[str] = set()
    shape_colors: Set[str] = set()
    background_colors: Set[str] = set()

    for p in doc.paragraphs:
        # Extract paragraph-level formatting
        docx_extract_deep_formatting(p._element, text_colors, shape_colors, fonts)
        
        for r in p.runs:
            try:
                if r.font and r.font.name:
                    fonts.add(r.font.name)
                elif p.style and p.style.font and p.style.font.name:
                    fonts.add(p.style.font.name)
            except Exception:
                pass
            try:
                c = r.font.color.rgb
                hexv = rgbcolor_to_hex(c)
                if hexv:
                    text_colors.add(hexv)
            except Exception:
                pass

    try:
        for table in doc.tables:
            docx_extract_deep_formatting(table._element, text_colors, shape_colors, fonts)
            for row in table.rows:
                for cell in row.cells:
                    docx_extract_deep_formatting(cell._element, text_colors, shape_colors, fonts)
                    for paragraph in cell.paragraphs:
                        docx_extract_deep_formatting(paragraph._element, text_colors, shape_colors, fonts)
                    shd_elems = cell._tc.xpath('.//w:shd')
                    for shd in shd_elems:
                        fill_val = shd.get(qn('w:fill'))
                        if fill_val and fill_val != "auto":
                            shape_colors.add("#" + fill_val.upper())
                    borders = cell._tc.xpath('.//w:tcBorders/*')
                    for b in borders:
                        col = b.get(qn('w:color'))
                        if col and col.lower() not in ("auto", "none"):
                            if len(col) == 3:
                                col = "".join([ch*2 for ch in col])
                            if len(col) == 6:
                                shape_colors.add("#" + col.upper())
    except Exception:
        pass

    images: List[Dict] = []
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes), 'r')
        for name in zf.namelist():
            if name.startswith("word/media/"):
                data = zf.read(name)
                images.append({"name": name.split("/")[-1], "path": name, "media_path": name, "uid": name, "bytes": data, "group": "Document"})
        for name in zf.namelist():
            if "embeddings" in name or "oleObject" in name:
                try:
                    data = zf.read(name)
                    images.append({"name": f"Embedded: {name.split('/')[-1]}", "path": name, "media_path": name, "uid": f"embed_{name}", "bytes": data, "group": "Embedded Objects"})
                except Exception:
                    pass
        zf.close()
    except Exception:
        pass

    return {"document": doc, "text_colors": sorted(list(text_colors)), "shape_colors": sorted(list(shape_colors)), "background_colors": sorted(list(background_colors)), "fonts": sorted(list(fonts)), "images": images}

def docx_apply_updates(extracted, color_map: Dict[str, str], font_map: Dict[str, str], image_replacements: Dict[str, bytes]) -> bytes:
    doc: DocxDocument = extracted["document"]
    for p in doc.paragraphs:
        for r in p.runs:
            try:
                current_font = r.font.name
                if current_font and current_font in font_map and font_map[current_font]:
                    r.font.name = font_map[current_font]
            except Exception:
                pass
            try:
                c = r.font.color.rgb
                curr_hex = rgbcolor_to_hex(c)
                if curr_hex and curr_hex in color_map and color_map[curr_hex]:
                    r.font.color.rgb = RGBColor.from_string(hex_no_hash(color_map[curr_hex]))
            except Exception:
                pass

    try:
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    shd_elems = cell._tc.xpath('.//w:shd')
                    for shd in shd_elems:
                        fill_val = shd.get(qn('w:fill'))
                        if fill_val and fill_val != "auto":
                            curr_hex = "#" + fill_val.upper()
                            if curr_hex in color_map and color_map[curr_hex]:
                                tcPr = cell._tc.get_or_add_tcPr()
                                new_shd = OxmlElement('w:shd')
                                new_shd.set(qn('w:fill'), hex_no_hash(color_map[curr_hex]))
                                try:
                                    for old in shd_elems:
                                        tcPr.remove(old)
                                except Exception:
                                    pass
                                tcPr.append(new_shd)
                    borders = cell._tc.xpath('.//w:tcBorders/*')
                    for b in borders:
                        col = b.get(qn('w:color'))
                        if col and col.lower() not in ("auto", "none"):
                            if len(col) == 3:
                                col_hex = "#" + "".join([ch*2 for ch in col]).upper()
                            else:
                                col_hex = "#" + col.upper()
                            if col_hex in color_map and color_map[col_hex]:
                                b.set(qn('w:color'), hex_no_hash(color_map[col_hex]))
    except Exception:
        pass

    out_buf = io.BytesIO()
    doc.save(out_buf)
    out_bytes = out_buf.getvalue()

    if image_replacements:
        try:
            in_zip = zipfile.ZipFile(io.BytesIO(out_bytes), 'r')
            out_mem = io.BytesIO()
            out_zip = zipfile.ZipFile(out_mem, 'w', zipfile.ZIP_DEFLATED)
            for name in in_zip.namelist():
                data = in_zip.read(name)
                if name in image_replacements and name.startswith("word/media/"):
                    target_ext = os.path.splitext(name)[1]
                    data = convert_image_bytes_to_ext(image_replacements[name], target_ext)
                out_zip.writestr(name, data)
            in_zip.close()
            out_zip.close()
            out_bytes = out_mem.getvalue()
        except Exception:
            pass
    return out_bytes


# PPTX helper functions
def pptx_get_background_image(slide_or_layout_or_master):
    try:
        bg_elm = slide_or_layout_or_master.background._element
        blips = bg_elm.xpath(".//a:blip")
        if blips:
            r_id = blips[0].get(pptx_qn('r:embed'))
            if r_id:
                part = slide_or_layout_or_master.part.related_parts[r_id]
                return part.blob, r_id, part
    except Exception:
        pass
    return None, None, None

def pptx_get_shape_fill_picture(shape):
    try:
        if hasattr(shape, "fill") and shape.fill and shape.fill.type == MSO_FILL.PICTURE:
            blip = shape.fill._fill.blipFill.blip
            if blip is not None:
                r_id = blip.get(pptx_qn('r:embed'))
                if r_id:
                    part = shape.part.related_parts[r_id]
                    return part.blob, r_id, part
    except Exception:
        pass
    return None, None, None

def pptx_get_line_hex(shape) -> Optional[str]:
    try:
        ln = shape.line
        if ln is None:
            return None
        try:
            if ln.fill is None or ln.fill.type in (None, MSO_FILL.BACKGROUND):
                return None
        except Exception:
            pass
        try:
            hexv = extract_color_from_pptx_color_obj(ln.color)
            if hexv:
                return hexv
        except Exception:
            pass
        try:
            if ln.fill and ln.fill.type == MSO_FILL.SOLID:
                hexv = extract_color_from_pptx_color_obj(ln.fill.fore_color)
                if hexv:
                    return hexv
        except Exception:
            pass
    except Exception:
        return None
    return None

def pptx_set_line_hex(shape, new_hex: str):
    try:
        ln = shape.line
        if ln is None:
            return
        try:
            if ln.fill is None or ln.fill.type in (None, MSO_FILL.BACKGROUND):
                return
        except Exception:
            pass
        updated = False
        try:
            curr = extract_color_from_pptx_color_obj(ln.color)
            if curr:
                ln.color.rgb = PPTX_RGBColor.from_string(hex_no_hash(new_hex))
                updated = True
        except Exception:
            pass
        if not updated:
            try:
                if ln.fill and ln.fill.type == MSO_FILL.SOLID:
                    ln.fill.fore_color.rgb = PPTX_RGBColor.from_string(hex_no_hash(new_hex))
            except Exception:
                pass
    except Exception:
        pass

def pptx_extract_text_formatting(text_frame, text_colors: Set[str], fonts: Set[str]):
    """Extract all text formatting including default colors and fonts"""
    try:
        # Process each paragraph
        for paragraph in text_frame.paragraphs:
            # Paragraph-level font
            try:
                pf = paragraph.font
                if pf:
                    if pf.name:
                        fonts.add(pf.name)
                    # Try to get color
                    if pf.color:
                        hexv = extract_color_from_pptx_color_obj(pf.color)
                        if hexv:
                            text_colors.add(hexv)
            except Exception:
                pass
            
            # Run-level formatting
            for run in paragraph.runs:
                try:
                    rf = run.font
                    if rf:
                        if rf.name:
                            fonts.add(rf.name)
                        # Try to get color
                        if rf.color:
                            hexv = extract_color_from_pptx_color_obj(rf.color)
                            if hexv:
                                text_colors.add(hexv)
                except Exception:
                    pass
                    
    except Exception:
        pass

def pptx_update_text_formatting(text_frame, color_map: Dict[str, str], font_map: Dict[str, str]):
    """Update all text formatting including paragraphs and runs"""
    try:
        for paragraph in text_frame.paragraphs:
            # Update paragraph-level font
            try:
                pf = paragraph.font
                if pf:
                    if pf.name and pf.name in font_map and font_map[pf.name]:
                        pf.name = font_map[pf.name]
                    if pf.color:
                        curr_hex = extract_color_from_pptx_color_obj(pf.color)
                        if curr_hex and curr_hex in color_map and color_map[curr_hex]:
                            pf.color.rgb = PPTX_RGBColor.from_string(hex_no_hash(color_map[curr_hex]))
            except Exception:
                pass
            
            # Update run-level formatting
            for run in paragraph.runs:
                try:
                    rf = run.font
                    if rf:
                        if rf.name and rf.name in font_map and font_map[rf.name]:
                            rf.name = font_map[rf.name]
                        if rf.color:
                            curr_hex = extract_color_from_pptx_color_obj(rf.color)
                            if curr_hex and curr_hex in color_map and color_map[curr_hex]:
                                rf.color.rgb = PPTX_RGBColor.from_string(hex_no_hash(color_map[curr_hex]))
                except Exception:
                    pass
    except Exception:
        pass

def pptx_process_shape_recursive(shape, slide_idx: int, path: str, text_colors: Set[str], shape_colors: Set[str], fonts: Set[str], images: List[Dict], depth: int = 0) -> None:
    """Recursively process all shape types including nested groups"""
    if depth > 10:
        return
    
    try:
        shape_type = shape.shape_type
        shape_name = getattr(shape, 'name', f'Shape_{path}')
        
        # CRITICAL: Process ALL shapes with text_frame, not just those with has_text_frame
        # This catches text boxes, shapes with text, and all other text containers
        try:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                pptx_extract_text_formatting(shape.text_frame, text_colors, fonts)
        except Exception:
            pass
        
        # Also check the traditional way
        if hasattr(shape, "has_text_frame"):
            try:
                if shape.has_text_frame and shape.text_frame is not None:
                    pptx_extract_text_formatting(shape.text_frame, text_colors, fonts)
            except Exception:
                pass

        # Process title and placeholders explicitly
        if hasattr(shape, "is_placeholder"):
            try:
                if shape.is_placeholder and hasattr(shape, "text_frame") and shape.text_frame:
                    pptx_extract_text_formatting(shape.text_frame, text_colors, fonts)
            except Exception:
                pass

        # Tables
        if hasattr(shape, "has_table") and shape.has_table:
            try:
                tbl = shape.table
                for row_idx, row in enumerate(tbl.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_path = f"{path}_tbl_r{row_idx}_c{col_idx}"
                        
                        # Cell text formatting
                        try:
                            tf = cell.text_frame
                            if tf:
                                pptx_extract_text_formatting(tf, text_colors, fonts)
                        except Exception:
                            pass
                        
                        try:
                            cell_fill = cell.fill
                            if cell_fill and cell_fill.type == MSO_FILL.SOLID:
                                hexv = extract_color_from_pptx_color_obj(cell_fill.fore_color)
                                if hexv:
                                    shape_colors.add(hexv)
                        except Exception:
                            pass
                        try:
                            if cell.fill and cell.fill.type == MSO_FILL.PICTURE:
                                blip = cell.fill._fill.blipFill.blip
                                if blip is not None:
                                    r_id = blip.get(pptx_qn('r:embed'))
                                    if r_id:
                                        part = cell._tc.part.related_parts[r_id]
                                        blob = part.blob
                                        images.append({"name": f"Slide {slide_idx+1} Table Cell Picture ({cell_path})", "bytes": blob, "uid": f"pptx_fill_{slide_idx}_{cell_path}", "group": f"Slide {slide_idx+1}", "kind": "cell_fill", "rel_id": r_id, "media_path": str(part.partname).lstrip("/") if part else None})
                        except Exception:
                            pass
            except Exception:
                pass

        # Fill colors
        if hasattr(shape, "fill"):
            try:
                fill = shape.fill
                if fill:
                    if fill.type == MSO_FILL.SOLID:
                        hexv = extract_color_from_pptx_color_obj(fill.fore_color)
                        if hexv:
                            shape_colors.add(hexv)
                    elif fill.type == MSO_FILL.PICTURE:
                        blob, r_id, part = pptx_get_shape_fill_picture(shape)
                        if blob:
                            images.append({"name": f"Slide {slide_idx+1} Shape Fill Picture ({shape_name})", "bytes": blob, "uid": f"pptx_fill_{slide_idx}_{path}", "group": f"Slide {slide_idx+1}", "kind": "shape_fill", "rel_id": r_id, "media_path": str(part.partname).lstrip("/") if part else None})
                    elif fill.type == MSO_FILL.GRADIENT:
                        try:
                            gradient = fill.gradient_stops
                            for stop in gradient:
                                hexv = extract_color_from_pptx_color_obj(stop.color)
                                if hexv:
                                    shape_colors.add(hexv)
                        except Exception:
                            pass
                    elif fill.type == MSO_FILL.PATTERNED:
                        try:
                            hexv = extract_color_from_pptx_color_obj(fill.fore_color)
                            if hexv:
                                shape_colors.add(hexv)
                            hexv = extract_color_from_pptx_color_obj(fill.back_color)
                            if hexv:
                                shape_colors.add(hexv)
                        except Exception:
                            pass
            except Exception:
                pass

        # Line colors
        try:
            line_hex = pptx_get_line_hex(shape)
            if line_hex:
                shape_colors.add(line_hex)
        except Exception:
            pass

        # Picture shapes
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                fname = getattr(shape.image, 'filename', 'image')
                r_id = shape._element.blipFill.blip.get(pptx_qn('r:embed'))
                part = shape.part.related_parts.get(r_id)
                media_path = str(part.partname).lstrip("/") if part is not None else None
                images.append({"name": f"Slide {slide_idx+1} Picture ({fname})", "bytes": blob, "uid": f"pptx_{slide_idx}_{path}", "group": f"Slide {slide_idx+1}", "kind": "shape_picture", "filename": fname, "media_path": media_path})
            except Exception:
                pass

        # Recursive group processing
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub_idx, sub_shape in enumerate(shape.shapes):
                    sub_path = f"{path}_g{sub_idx}"
                    pptx_process_shape_recursive(sub_shape, slide_idx, sub_path, text_colors, shape_colors, fonts, images, depth + 1)
            except Exception:
                pass
                
    except Exception:
        pass

def pptx_update_shape_recursive(shape, slide_idx: int, path: str, color_map: Dict[str, str], font_map: Dict[str, str], image_replacements: Dict[str, bytes], depth: int = 0) -> None:
    """Recursively update all shape types including nested groups"""
    if depth > 10:
        return
    
    try:
        shape_type = shape.shape_type
        
        # CRITICAL: Update ALL shapes with text_frame, not just those with has_text_frame
        # This catches text boxes, shapes with text, and all other text containers
        try:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                pptx_update_text_formatting(shape.text_frame, color_map, font_map)
        except Exception:
            pass
        
        # Also check the traditional way
        if hasattr(shape, "has_text_frame"):
            try:
                if shape.has_text_frame and shape.text_frame is not None:
                    pptx_update_text_formatting(shape.text_frame, color_map, font_map)
            except Exception:
                pass

        # Update title and placeholders explicitly
        if hasattr(shape, "is_placeholder"):
            try:
                if shape.is_placeholder and hasattr(shape, "text_frame") and shape.text_frame:
                    pptx_update_text_formatting(shape.text_frame, color_map, font_map)
            except Exception:
                pass

        # Update tables
        if hasattr(shape, "has_table") and shape.has_table:
            try:
                tbl = shape.table
                for row_idx, row in enumerate(tbl.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_path = f"{path}_tbl_r{row_idx}_c{col_idx}"
                        
                        try:
                            tf = cell.text_frame
                            if tf:
                                pptx_update_text_formatting(tf, color_map, font_map)
                        except Exception:
                            pass
                        
                        try:
                            cell_fill = cell.fill
                            if cell_fill and cell_fill.type == MSO_FILL.SOLID:
                                curr_hex = extract_color_from_pptx_color_obj(cell_fill.fore_color)
                                if curr_hex and curr_hex in color_map and color_map[curr_hex]:
                                    cell_fill.solid()
                                    cell_fill.fore_color.rgb = PPTX_RGBColor.from_string(hex_no_hash(color_map[curr_hex]))
                        except Exception:
                            pass
                        try:
                            uid_cell = f"pptx_fill_{slide_idx}_{cell_path}"
                            if uid_cell in image_replacements and cell.fill and cell.fill.type == MSO_FILL.PICTURE:
                                blip = cell.fill._fill.blipFill.blip
                                if blip is not None:
                                    r_id = blip.get(pptx_qn('r:embed'))
                                    if r_id:
                                        part = cell._tc.part.related_parts[r_id]
                                        ext = os.path.splitext(str(part.partname))[-1] if hasattr(part, "partname") else ".png"
                                        part.blob = convert_image_bytes_to_ext(image_replacements[uid_cell], ext)
                        except Exception:
                            pass
            except Exception:
                pass

        # Update fills
        if hasattr(shape, "fill"):
            try:
                fill = shape.fill
                if fill and fill.type == MSO_FILL.SOLID:
                    curr_hex = extract_color_from_pptx_color_obj(fill.fore_color)
                    if curr_hex and curr_hex in color_map and color_map[curr_hex]:
                        fill.solid()
                        fill.fore_color.rgb = PPTX_RGBColor.from_string(hex_no_hash(color_map[curr_hex]))
            except Exception:
                pass

        # Update lines
        try:
            line_hex = pptx_get_line_hex(shape)
            if line_hex and line_hex in color_map:
                pptx_set_line_hex(shape, color_map[line_hex])
        except Exception:
            pass

        # Update pictures
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                uid = f"pptx_{slide_idx}_{path}"
                if uid in image_replacements:
                    r_id = shape._element.blipFill.blip.get(pptx_qn('r:embed'))
                    image_part = shape.part.related_parts[r_id]
                    ext = os.path.splitext(getattr(shape.image, "filename", "image.png"))[-1] or ".png"
                    image_part.blob = convert_image_bytes_to_ext(image_replacements[uid], ext)
            except Exception:
                pass

        # Update picture fills
        try:
            uid_fill = f"pptx_fill_{slide_idx}_{path}"
            if uid_fill in image_replacements and hasattr(shape, "fill") and shape.fill and shape.fill.type == MSO_FILL.PICTURE:
                blip = shape.fill._fill.blipFill.blip
                if blip is not None:
                    r_id = blip.get(pptx_qn('r:embed'))
                    if r_id:
                        part = shape.part.related_parts[r_id]
                        ext = os.path.splitext(str(part.partname))[-1] if hasattr(part, "partname") else ".png"
                        part.blob = convert_image_bytes_to_ext(image_replacements[uid_fill], ext)
        except Exception:
            pass

        # Recursive group processing
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub_idx, sub_shape in enumerate(shape.shapes):
                    sub_path = f"{path}_g{sub_idx}"
                    pptx_update_shape_recursive(sub_shape, slide_idx, sub_path, color_map, font_map, image_replacements, depth + 1)
            except Exception:
                pass
                
    except Exception:
        pass

def pptx_compose_slide_preview(prs: Presentation, slide, width_px: int = 900) -> Optional[bytes]:
    if not PIL_AVAILABLE:
        return None
    try:
        slide_w = int(prs.slide_width)
        slide_h = int(prs.slide_height)
        ratio = width_px / float(slide_w)
        height_px = int(slide_h * ratio)
        canvas = PILImage.new("RGBA", (width_px, height_px), (255, 255, 255, 255))

        bg_blob, _, _ = pptx_get_background_image(slide)
        if bg_blob:
            bg_img = PILImage.open(io.BytesIO(bg_blob)).convert("RGBA")
            bg_img = bg_img.resize((width_px, height_px))
            canvas.paste(bg_img, (0, 0), bg_img if bg_img.mode == "RGBA" else None)

        def paste_pictures(shape, offset_left_emu=0, offset_top_emu=0):
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blob = shape.image.blob
                    pic = PILImage.open(io.BytesIO(blob)).convert("RGBA")
                    left = int((int(shape.left) + offset_left_emu) * ratio)
                    top = int((int(shape.top) + offset_top_emu) * ratio)
                    w = max(1, int(int(shape.width) * ratio))
                    h = max(1, int(int(shape.height) * ratio))
                    pic = pic.resize((w, h))
                    canvas.paste(pic, (left, top), pic)
            except Exception:
                pass
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub in shape.shapes:
                        paste_pictures(sub, offset_left_emu + int(shape.left), offset_top_emu + int(shape.top))
            except Exception:
                pass

        for shape in slide.shapes:
            paste_pictures(shape)

        out = io.BytesIO()
        canvas.convert("RGB").save(out, format="JPEG", quality=85)
        return out.getvalue()
    except Exception:
        return None

def pptx_extract_theme_images(file_bytes: bytes):
    themes = []
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes), 'r')
        theme_paths = [n for n in zf.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")]
        for tpath in theme_paths:
            images = []
            rels_path = "ppt/theme/_rels/" + tpath.split("/")[-1] + ".rels"
            if rels_path in zf.namelist():
                rels_xml = zf.read(rels_path)
                try:
                    rels_root = ET.fromstring(rels_xml)
                    for rel in rels_root.findall(".//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"):
                        rId = rel.attrib.get("Id")
                        target = rel.attrib.get("Target", "")
                        if "media/" in target:
                            media_path = normalize_zip_path("ppt/theme", target)
                            if media_path in zf.namelist():
                                img_bytes = zf.read(media_path)
                                uid = f"pptx_theme_img_{tpath.split('/')[-1]}_{rId}"
                                images.append({"uid": uid, "media_path": media_path, "bytes": img_bytes, "rid": rId, "theme_path": tpath, "name": media_path.split("/")[-1], "group": f"Theme {tpath.split('/')[-1]}"})
                except Exception:
                    pass
            if images:
                themes.append({"theme_path": tpath, "images": images})
        zf.close()
    except Exception:
        pass
    return {"themes": themes}

def pptx_list_all_media(file_bytes: bytes) -> List[Dict]:
    items = []
    try:
        zf = zipfile.ZipFile(io.BytesIO(file_bytes), 'r')
        for name in zf.namelist():
            if name.startswith("ppt/media/"):
                try:
                    data = zf.read(name)
                    items.append({"name": name.split("/")[-1], "uid": name, "media_path": name, "bytes": data, "group": "All Media"})
                except Exception:
                    continue
        zf.close()
    except Exception:
        pass
    return items

def pptx_extract(file_bytes: bytes):
    if not PPTX_AVAILABLE:
        return None

    buf = io.BytesIO(file_bytes)
    prs = Presentation(buf)

    text_colors: Set[str] = set()
    shape_colors: Set[str] = set()
    background_colors: Set[str] = set()
    fonts: Set[str] = set()
    images: List[Dict] = []
    slide_previews: Dict[str, bytes] = {}

    # Master background
    try:
        master = getattr(prs, "slide_master", None)
        if master is not None:
            blob, r_id, part = pptx_get_background_image(master)
            if blob:
                images.append({"name": "Slide Master Background", "bytes": blob, "uid": f"pptx_master_bg_0", "group": "Master", "kind": "master_bg", "rel_id": r_id, "media_path": str(part.partname).lstrip("/") if part else None})
    except Exception:
        pass

    # Layout backgrounds
    for layout_idx, layout in enumerate(prs.slide_layouts):
        blob, r_id, part = pptx_get_background_image(layout)
        if blob:
            images.append({"name": f"Layout {layout_idx+1} Background", "bytes": blob, "uid": f"pptx_layout_bg_{layout_idx}", "group": f"Layout {layout_idx+1}", "kind": "layout_bg", "rel_id": r_id, "media_path": str(part.partname).lstrip("/") if part else None})

    # Slides
    for slide_idx, slide in enumerate(prs.slides):
        try:
            fill = slide.background.fill
            if fill and fill.type == MSO_FILL.SOLID:
                hexv = extract_color_from_pptx_color_obj(fill.fore_color)
                if hexv:
                    background_colors.add(hexv)
        except Exception:
            pass

        blob, r_id, part = pptx_get_background_image(slide)
        if blob:
            images.append({"name": f"Slide {slide_idx+1} Background", "bytes": blob, "uid": f"pptx_slide_bg_{slide_idx}", "group": f"Slide {slide_idx+1}", "kind": "slide_bg", "rel_id": r_id, "media_path": str(part.partname).lstrip("/") if part else None})

        for shape_idx, shape in enumerate(slide.shapes):
            path = str(shape_idx)
            pptx_process_shape_recursive(shape, slide_idx, path, text_colors, shape_colors, fonts, images, depth=0)

        preview_bytes = pptx_compose_slide_preview(prs, slide, width_px=900) if PIL_AVAILABLE else None
        if preview_bytes:
            slide_previews[f"Slide {slide_idx+1}"] = preview_bytes

    theme_images_info = pptx_extract_theme_images(file_bytes)
    all_media = pptx_list_all_media(file_bytes)
    images.extend(all_media)

    return {"presentation": prs, "text_colors": sorted(list(text_colors)), "shape_colors": sorted(list(shape_colors)), "background_colors": sorted(list(background_colors)), "fonts": sorted(list(fonts)), "images": images, "slide_previews": slide_previews, "theme_images_info": theme_images_info}

def zip_replace_media(base_bytes: bytes, replacements: Dict[str, bytes]) -> bytes:
    if not replacements:
        return base_bytes
    try:
        in_zip = zipfile.ZipFile(io.BytesIO(base_bytes), 'r')
        out_mem = io.BytesIO()
        out_zip = zipfile.ZipFile(out_mem, 'w', zipfile.ZIP_DEFLATED)
        for name in in_zip.namelist():
            data = in_zip.read(name)
            if name in replacements:
                data = replacements[name]
            out_zip.writestr(name, data)
        in_zip.close()
        out_zip.close()
        return out_mem.getvalue()
    except Exception:
        return base_bytes

def pptx_apply_updates(extracted, color_map: Dict[str, str], font_map: Dict[str, str], image_replacements: Dict[str, bytes], theme_image_replacements: Dict[str, bytes]) -> bytes:
    prs: Presentation = extracted["presentation"]

    uid_to_media: Dict[str, Optional[str]] = {}
    for img in extracted.get("images", []):
        uid_to_media[img.get("uid")] = img.get("media_path")

    for slide_idx, slide in enumerate(prs.slides):
        try:
            fill = slide.background.fill
            if fill and fill.type == MSO_FILL.SOLID:
                curr_hex = extract_color_from_pptx_color_obj(fill.fore_color)
                if curr_hex and curr_hex in color_map and color_map[curr_hex]:
                    fill.solid()
                    fill.fore_color.rgb = PPTX_RGBColor.from_string(hex_no_hash(color_map[curr_hex]))
        except Exception:
            pass

        try:
            uid_bg = f"pptx_slide_bg_{slide_idx}"
            if uid_bg in image_replacements:
                bg_elm = slide.background._element
                blips = bg_elm.xpath(".//a:blip")
                if blips:
                    r_id = blips[0].get(pptx_qn('r:embed'))
                    if r_id:
                        part = slide.part.related_parts[r_id]
                        ext = os.path.splitext(str(part.partname))[-1] if hasattr(part, "partname") else ".png"
                        part.blob = convert_image_bytes_to_ext(image_replacements[uid_bg], ext)
        except Exception:
            pass

        for shape_idx, shape in enumerate(slide.shapes):
            path = str(shape_idx)
            pptx_update_shape_recursive(shape, slide_idx, path, color_map, font_map, image_replacements, depth=0)

    out_buf = io.BytesIO()
    prs.save(out_buf)
    updated_pptx = out_buf .getvalue()

    zip_media_repls: Dict[str, bytes] = {}
    for uid, media_path in uid_to_media.items():
        if uid in image_replacements and media_path:
            target_ext = os.path.splitext(media_path)[1] or ".png"
            zip_media_repls[media_path] = convert_image_bytes_to_ext(image_replacements[uid], target_ext)
    for uid, data in image_replacements.items():
        if uid.startswith("ppt/media/"):
            target_ext = os.path.splitext(uid)[1] or ".png"
            zip_media_repls[uid] = convert_image_bytes_to_ext(data, target_ext)
    for media_path, data in theme_image_replacements.items():
        target_ext = os.path.splitext(media_path)[1] or ".png"
        zip_media_repls[media_path] = convert_image_bytes_to_ext(data, target_ext)

    if zip_media_repls:
        updated_pptx = zip_replace_media(updated_pptx, zip_media_repls)

    return updated_pptx


# XLSX functions
def xlsx_extract(file_bytes: bytes):
    if not OPENPYXL_AVAILABLE:
        return None

    buf = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(buf, data_only=True)

    text_colors: Set[str] = set()
    shape_colors: Set[str] = set()
    background_colors: Set[str] = set()
    fonts: Set[str] = set()
    images: List[Dict] = []

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                try:
                    if cell.font and cell.font.name:
                        fonts.add(cell.font.name)
                except Exception:
                    pass
                try:
                    hexv = openpyxl_color_to_hex(cell.font.color)
                    if hexv:
                        text_colors.add(hexv)
                except Exception:
                    pass
                try:
                    fill = cell.fill
                    if fill and fill.patternType == "solid":
                        hexv = openpyxl_color_to_hex(fill.fgColor)
                        if hexv:
                            shape_colors.add(hexv)
                            background_colors.add(hexv)
                except Exception:
                    pass
                try:
                    b = cell.border
                    if b:
                        for side_name in ["left", "right", "top", "bottom"]:
                            side = getattr(b, side_name)
                            if side and side.color:
                                hexv = openpyxl_color_to_hex(side.color)
                                if hexv:
                                    shape_colors.add(hexv)
                except Exception:
                    pass

        try:
            ws_images = getattr(ws, "_images", [])
            for idx, img in enumerate(ws_images):
                uid = f"xlsx_{ws.title}_{idx}"
                preview_bytes = None
                if PIL_AVAILABLE:
                    try:
                        pil_img = getattr(img, "_data", None)
                        if pil_img is not None:
                            out = io.BytesIO()
                            pil_img.save(out, format="PNG")
                            preview_bytes = out.getvalue()
                    except Exception:
                        preview_bytes = None
                images.append({"name": f"{ws.title} Image {idx+1}", "bytes": preview_bytes, "uid": uid, "group": f"Sheet {ws.title}", "ws_title": ws.title, "index": idx})
        except Exception:
            pass

    return {"workbook": wb, "text_colors": sorted(list(text_colors)), "shape_colors": sorted(list(shape_colors)), "background_colors": sorted(list(background_colors)), "fonts": sorted(list(fonts)), "images": images}

def xlsx_apply_updates(extracted, color_map: Dict[str, str], font_map: Dict[str, str], image_replacements: Dict[str, bytes]) -> bytes:
    wb = extracted["workbook"]

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                try:
                    curr_font = cell.font
                    new_name = None
                    if curr_font and curr_font.name and curr_font.name in font_map and font_map[curr_font.name]:
                        new_name = font_map[curr_font.name]
                    curr_color_hex = openpyxl_color_to_hex(curr_font.color) if curr_font and curr_font.color else None
                    new_color_hex = None
                    if curr_color_hex and curr_color_hex in color_map and color_map[curr_color_hex]:
                        new_color_hex = color_map[curr_color_hex]
                    if new_name or new_color_hex:
                        kwargs = {}
                        if new_name:
                            kwargs["name"] = new_name
                        if new_color_hex:
                            kwargs["color"] = Color(rgb="FF" + hex_no_hash(new_color_hex))
                        cell.font = Font(name=kwargs.get("name", curr_font.name), size=curr_font.size, bold=curr_font.bold, italic=curr_font.italic, vertAlign=curr_font.vertAlign, underline=curr_font.underline, strike=curr_font.strike, color=kwargs.get("color", curr_font.color), shadow=curr_font.shadow, scheme=curr_font.scheme, charset=curr_font.charset, outline=curr_font.outline, condense=curr_font.condense, extend=curr_font.extend)
                except Exception:
                    pass

                try:
                    fill = cell.fill
                    if fill and fill.patternType == "solid":
                        curr_fill_hex = openpyxl_color_to_hex(fill.fgColor)
                        if curr_fill_hex and curr_fill_hex in color_map and color_map[curr_fill_hex]:
                            new_hex = color_map[curr_fill_hex]
                            cell.fill = PatternFill(fill_type="solid", fgColor=Color(rgb="FF" + hex_no_hash(new_hex)))
                except Exception:
                    pass

                try:
                    b = cell.border
                    if b:
                        sides = {}
                        for side_name in ["left", "right", "top", "bottom"]:
                            side = getattr(b, side_name)
                            if side:
                                hexv = openpyxl_color_to_hex(side.color) if side.color else None
                                if hexv and hexv in color_map and color_map[hexv]:
                                    new_side = Side(style=side.style, color=Color(rgb="FF" + hex_no_hash(color_map[hexv])))
                                else:
                                    new_side = side
                                sides[side_name] = new_side
                        cell.border = Border(left=sides.get("left", b.left), right=sides.get("right", b.right), top=sides.get("top", b.top), bottom=sides.get("bottom", b.bottom), diagonal=b.diagonal, diagonalDown=b.diagonalDown, diagonalUp=b.diagonalUp, outline=b.outline, vertical=b.vertical, horizontal=b.horizontal)
                except Exception:
                    pass

    if image_replacements and PIL_AVAILABLE:
        try:
            for ws in wb.worksheets:
                ws_images = getattr(ws, "_images", [])
                new_images = []
                for idx, img in enumerate(ws_images):
                    anchor = getattr(img, "anchor", None)
                    uid = f"xlsx_{ws.title}_{idx}"
                    if uid in image_replacements:
                        rep_bytes = image_replacements[uid]
                        rep_buf = io.BytesIO(rep_bytes)
                        pil_img = PILImage.open(rep_buf).convert("RGBA")
                        xl_img = XLImage(pil_img)
                        if anchor:
                            xl_img.anchor = anchor
                        new_images.append(xl_img)
                    else:
                        new_images.append(img)
                ws._images = new_images
        except Exception:
            pass

    out_buf = io.BytesIO()
    wb.save(out_buf)
    return out_buf.getvalue()


# PDF
def pdf_preview(file_bytes: bytes):
    st.warning("PDF preview and in-place full rebranding are limited in this tool. You can download and review the uploaded PDF below.")
    st.download_button("Download uploaded PDF", data=file_bytes, file_name="uploaded.pdf")


# Main UI
st.markdown('<div class="pwc-header"><h2>PwC Rebranding Tool</h2><div class="pwc-subtle">Upload a document and guide the rebranding of colors, fonts, and images.</div></div>', unsafe_allow_html=True)

uploaded = st.file_uploader("Upload your document", type=["docx", "pptx", "xlsx", "pdf"])

if uploaded is None:
    st.info("Please upload a document to begin.")
    st.stop()

file_bytes = uploaded.read()
file_type = infer_file_type(uploaded.name)

if file_type is None:
    st.error("Unsupported file type. Please upload a .docx, .pptx, .xlsx, or .pdf file.")
    st.stop()

st.markdown('<div class="pwc-card"><div class="pwc-section-title">Uploaded Document</div>', unsafe_allow_html=True)
st.write(f"File name: {uploaded.name}")
st.write(f"Type: {file_type.upper()}")

if file_type == "pdf":
    pdf_preview(file_bytes)
else:
    st.markdown('<div class="pwc-hint">A full visual rendering is not always available, but a structured preview is provided below.</div>', unsafe_allow_html=True)

st.markdown("</div>", unsafe_allow_html=True)

# Extract metadata
extracted = None
if file_type == "docx":
    if not DOCX_AVAILABLE:
        st.error("python-docx not installed. Please install with: pip install python-docx")
        st.stop()
    extracted = docx_extract(file_bytes)
elif file_type == "pptx":
    if not PPTX_AVAILABLE:
        st.error("python-pptx not installed. Please install with: pip install python-pptx")
        st.stop()
    extracted = pptx_extract(file_bytes)
elif file_type == "xlsx":
    if not OPENPYXL_AVAILABLE:
        st.error("openpyxl not installed. Please install with: pip install openpyxl")
        st.stop()
    extracted = xlsx_extract(file_bytes)
elif file_type == "pdf":
    extracted = None

if file_type != "pdf" and extracted is None:
    st.error("Failed to parse the uploaded document.")
    st.stop()

# Step 1: Colors & Fonts
if file_type != "pdf":
    col1, col2 = st.columns(2)
    with col1:
        st.markdown('<div class="pwc-card"><div class="pwc-section-title">Step 1: Colors</div>', unsafe_allow_html=True)
        txt_colors = extracted["text_colors"]
        shp_colors = extracted["shape_colors"]
        bg_colors = extracted["background_colors"]

        st.write("Text colors (includes all text in shapes, text boxes, titles, subtitles, bullets, and numbering):")
        if txt_colors:
            for c in txt_colors:
                st.color_picker(f"Change text color {c}", value=c, key=f"text_color_{safe_key(c)}")
        else:
            st.write("- None detected")

        st.write("Shapes/format colors (incl. borders):")
        if shp_colors:
            for c in shp_colors:
                st.color_picker(f"Change shape/border color {c}", value=c, key=f"shape_color_{safe_key(c)}")
        else:
            st.write("- None detected")

        st.write("Background colors:")
        if bg_colors:
            for c in bg_colors:
                st.color_picker(f"Change background color {c}", value=c, key=f"bg_color_{safe_key(c)}")
        else:
            st.write("- None detected or not supported for this file type")
        st.markdown("</div>", unsafe_allow_html=True)

    with col2:
        st.markdown('<div class="pwc-card"><div class="pwc-section-title">Step 1: Fonts</div>', unsafe_allow_html=True)
        fonts = extracted["fonts"]
        if fonts:
            for f in fonts:
                st.text_input(f"Change font '{f}' to:", value=f, key=f"font_map_{safe_key(f)}")
        else:
            st.write("- None detected")
        st.markdown("</div>", unsafe_allow_html=True)

# Step 2: Images
st.markdown('<div class="pwc-card"><div class="pwc-section-title">Step 2: Images, Pictures, or Logos</div>', unsafe_allow_html=True)

left_control, right_control = st.columns([3, 1])
with right_control:
    skip_images = st.checkbox("Skip images", value=False)
    thumb_width = st.slider("Thumbnail width (px)", min_value=80, max_value=300, value=140)
    if st.button("Clear all replacements"):
        clear_all_replacements()
        st.success("Cleared all image replacements.")

image_replacements_ss, theme_image_replacements_ss = get_persisted_image_replacements()

if file_type == "pdf":
    st.info("PDF image replacement is not supported in this version.")
else:
    if skip_images:
        st.info("Skipping image review and replacements.")
    else:
        images = extracted["images"]
        if file_type == "pptx":
            slide_previews = extracted.get("slide_previews", {})
            slides: Dict[str, List[Dict]] = {}
            others: Dict[str, List[Dict]] = {}
            all_media_items: List[Dict] = []
            for img in images:
                grp = img.get("group") or "Other"
                if grp.startswith("Slide "):
                    slides.setdefault(grp, []).append(img)
                elif grp == "All Media":
                    all_media_items.append(img)
                else:
                    others.setdefault(grp, []).append(img)

            st.caption("Preview shows the full slide BEFORE changes on the left; replace images on the right.")

            for grp_name in sorted(slides.keys(), key=lambda x: int(x.split(" ")[1])):
                grp_imgs = slides[grp_name]
                with st.expander(f"{grp_name}", expanded=False):
                    left, right = st.columns([2, 3])
                    with left:
                        st.markdown("Before preview")
                        preview = slide_previews.get(grp_name)
                        if preview:
                            st.image(preview, use_container_width=True)
                        else:
                            bg = next((i for i in grp_imgs if i.get("kind") == "slide_bg" and i.get("bytes")), None)
                            if bg:
                                st.image(bg["bytes"], use_container_width=True)
                            else:
                                st.info("No slide preview available.")
                        st.caption("Slide images (mini thumbnails):")
                        mini = [i for i in grp_imgs if i.get("kind") in ("shape_picture", "shape_fill", "cell_fill") and i.get("bytes")]
                        if mini:
                            st.markdown('<div class="thumb-row">', unsafe_allow_html=True)
                            for m in mini:
                                st.markdown('<div class="thumb-item">', unsafe_allow_html=True)
                                try:
                                    st.image(m["bytes"], width=100)
                                except Exception:
                                    st.write("(Unavailable)")
                                st.markdown('</div>', unsafe_allow_html=True)
                            st.markdown('</div>', unsafe_allow_html=True)
                        else:
                            st.write("No picture shapes detected.")
                    with right:
                        st.markdown("Image replacements")
                        for idx, img in enumerate(grp_imgs):
                            uid = img.get("uid") or f"{file_type}_img_{grp_name}_{idx}"
                            st.write(img.get("name", "Unnamed"))
                            if img.get("bytes"):
                                try:
                                    st.image(img["bytes"], width=thumb_width)
                                except Exception:
                                    st.write("(Preview unavailable)")
                            else:
                                st.write("(Preview unavailable)")
                            rep = st.file_uploader("Replace image (optional)", type=["png", "jpg", "jpeg", "gif"], key=f"replace_{safe_key(uid)}")
                            if rep is not None:
                                persist_image_replacement(uid, rep.read())

            for grp_name, grp_imgs in others.items():
                with st.expander(f"{grp_name} images", expanded=False):
                    for idx, img in enumerate(grp_imgs):
                        uid = img.get("uid") or f"{file_type}_img_{grp_name}_{idx}"
                        st.write(img.get("name", "Unnamed"))
                        if img.get("bytes"):
                            try:
                                st.image(img["bytes"], width=thumb_width)
                            except Exception:
                                st.write("(Preview unavailable)")
                        else:
                            st.write("(Preview unavailable)")
                        rep = st.file_uploader("Replace image (optional)", type=["png", "jpg", "jpeg", "gif"], key=f"replace_{safe_key(uid)}")
                        if rep is not None:
                            persist_image_replacement(uid, rep.read())

            theme_images_info = extracted.get("theme_images_info", {})
            if theme_images_info.get("themes"):
                with st.expander("Theme images (background assets)", expanded=False):
                    for theme in theme_images_info["themes"]:
                        for ti in theme["images"]:
                            st.write(f"{ti['name']} ({ti['media_path']})")
                            try:
                                st.image(ti["bytes"], width=thumb_width)
                            except Exception:
                                st.write("(Preview unavailable)")
                            rep = st.file_uploader("Replace theme image (optional)", type=["png", "jpg", "jpeg", "gif"], key=f"replace_theme_{safe_key(ti['uid'])}")
                            if rep is not None:
                                persist_theme_image_replacement(ti["media_path"], rep.read())

            if all_media_items:
                with st.expander("All Media (ppt/media) - fallback", expanded=False):
                    for idx, img in enumerate(all_media_items):
                        uid = img.get("media_path")
                        st.write(f"{img.get('name', 'media')} ({uid})")
                        if img.get("bytes"):
                            try:
                                st.image(img["bytes"], width=thumb_width)
                            except Exception:
                                st.write("(Preview unavailable)")
                        rep = st.file_uploader("Replace media (optional)", type=["png", "jpg", "jpeg", "gif"], key=f"replace_{safe_key(uid)}")
                        if rep is not None:
                            persist_image_replacement(uid, rep.read())

        else:
            groups: Dict[str, List[Dict]] = {}
            for img in images:
                grp = img.get("group") or "Other"
                groups.setdefault(grp, []).append(img)
            for grp_name, grp_imgs in groups.items():
                with st.expander(f"{grp_name} images", expanded=False):
                    for idx, img in enumerate(grp_imgs):
                        uid = img.get("uid") or f"{file_type}_img_{grp_name}_{idx}"
                        st.write(img.get("name", "Unnamed"))
                        if img.get("bytes"):
                            try:
                                st.image(img["bytes"], width=thumb_width)
                            except Exception:
                                st.write("(Preview unavailable)")
                        else:
                            st.write("(Preview unavailable)")
                        rep = st.file_uploader("Replace image (optional)", type=["png", "jpg", "jpeg", "gif"], key=f"replace_{safe_key(uid)}")
                        if rep is not None:
                            persist_image_replacement(uid, rep.read())

st.markdown("</div>", unsafe_allow_html=True)

# Apply rebranding
st.markdown('<div class="pwc-card"><div class="pwc-section-title">Apply Rebranding</div>', unsafe_allow_html=True)
apply_btn = st.button("Apply color/font changes and image replacements")

updated_bytes = None
updated_name = None

if apply_btn:
    if file_type == "pdf":
        st.warning("Rebranding changes are not applied to PDF files in this version.")
        updated_bytes = file_bytes
        updated_name = uploaded.name
    else:
        color_map: Dict[str, str] = {}
        for c in extracted["text_colors"]:
            new_c = st.session_state.get(f"text_color_{safe_key(c)}", c)
            if new_c and new_c != c:
                color_map[c] = new_c
        for c in extracted["shape_colors"]:
            new_c = st.session_state.get(f"shape_color_{safe_key(c)}", c)
            if new_c and new_c != c:
                color_map[c] = new_c
        for c in extracted["background_colors"]:
            new_c = st.session_state.get(f"bg_color_{safe_key(c)}", c)
            if new_c and new_c != c:
                color_map[c] = new_c

        font_map: Dict[str, str] = {}
        for f in extracted["fonts"]:
            new_f = st.session_state.get(f"font_map_{safe_key(f)}", f)
            if new_f and new_f != f:
                font_map[f] = new_f

        image_replacements, theme_image_replacements = get_persisted_image_replacements()

        try:
            if file_type == "docx":
                updated_bytes = docx_apply_updates(extracted, color_map, font_map, image_replacements)
                updated_name = uploaded.name.replace(".docx", "_rebranded.docx")
            elif file_type == "pptx":
                updated_bytes = pptx_apply_updates(extracted, color_map, font_map, image_replacements, theme_image_replacements)
                updated_name = uploaded.name.replace(".pptx", "_rebranded.pptx")
            elif file_type == "xlsx":
                updated_bytes = xlsx_apply_updates(extracted, color_map, font_map, image_replacements)
                updated_name = uploaded.name.replace(".xlsx", "_rebranded.xlsx")
        except Exception as e:
            st.error(f"Failed to apply updates: {e}")

    if updated_bytes:
        st.success("Rebranding applied successfully.")
        if file_type == "docx":
            st.markdown('<div class="pwc-hint">Please download and review the updated Word document.</div>', unsafe_allow_html=True)
        elif file_type == "pptx":
            st.markdown('<div class="pwc-hint">Slides updated. Download and review in PowerPoint.</div>', unsafe_allow_html=True)
        elif file_type == "xlsx":
            st.markdown('<div class="pwc-hint">Sheets updated. Download and review in Excel.</div>', unsafe_allow_html=True)
        elif file_type == "pdf":
            st.markdown('<div class="pwc-hint">No changes applied to PDF; download the original.</div>', unsafe_allow_html=True)

        st.download_button("Download rebranded document", data=updated_bytes, file_name=updated_name or uploaded.name)

st.markdown("</div>", unsafe_allow_html=True)

st.markdown(f"<div class='custom-footer'>{FOOTER_TEXT}</div>", unsafe_allow_html=True)

